import os, sys
import pandas as pd
import torch
from PIL import Image
import numpy as np
import torch.nn.functional as F
import cv2
from moviepy.editor import ImageSequenceClip
import matplotlib.pyplot as plt
from pptx import Presentation
from pptx.util import Inches, Pt
from math import isnan
from utils import load_image

def apply_registration(image_tensor, seg_tensor, grid_path):

    # handle unregistered images
    if isinstance(grid_path, float) and isnan(grid_path):
        return torch.zeros_like(image_tensor).squeeze(0), torch.zeros_like(seg_tensor).squeeze(0)

    # Load the sampling grid
    grid = torch.load(grid_path)[1]
    
    # Apply the sampling grid
    registered_image = F.grid_sample(image_tensor, grid)
    registered_seg = F.grid_sample(seg_tensor, grid)
    
    return registered_image.squeeze(0), registered_seg.squeeze(0)

def tensor_to_numpy(tensor):
    return tensor.permute(1, 2, 0).numpy()

def draw_contours(image_tensor, seg_tensor):
    image_np = tensor_to_numpy(image_tensor)
    seg_np = tensor_to_numpy(seg_tensor)

    image_np = (image_np * 255).astype(np.uint8)
    seg_np = (seg_np * 255).astype(np.uint8)

    if image_np.ndim == 2:
        image_np = cv2.cvtColor(image_np, cv2.COLOR_GRAY2BGR)
    
    contours, _ = cv2.findContours((seg_np * 255).astype(np.uint8), cv2.RETR_EXTERNAL, cv2.CHAIN_APPROX_SIMPLE)
    image_with_contours = cv2.drawContours(image_np.copy(), contours, -1, (0, 255, 0), 1)
    
    return torch.tensor(image_with_contours).permute(2, 0, 1)

def create_videos_and_plots(df, save_to, video_wcontours_name, video_wocontours_name, plot_name, thumbnails_wcontours_name, thumbnails_wocontours_name):

    # make videos and plots folder
    os.makedirs(os.path.join(save_to, 'metadata', 'videos_wcontours'), exist_ok=True)
    os.makedirs(os.path.join(save_to, 'metadata', 'videos_wocontours'), exist_ok=True)
    os.makedirs(os.path.join(save_to, 'metadata', 'plots'), exist_ok=True)
    os.makedirs(os.path.join(save_to, 'metadata', 'thumbnails_wcontours'), exist_ok=True)
    os.makedirs(os.path.join(save_to, 'metadata', 'thumbnails_wocontours'), exist_ok=True)

    # Initialize lists for frames and areas
    frames_with_contours = []
    frames_without_contours = []
    areas_ai = []
    areas_manual = []
    timepoints = []

    for i, row in df.iterrows():
        image_tensor = load_image(row.file_path_coris)
        seg_tensor = load_image(row.file_path_ga_segmentation)
        if i > 0:
            registered_image, registered_seg = apply_registration(image_tensor, seg_tensor, row.params)
            registered_seg = (registered_seg > 0.5).int()
        else:
            registered_image, registered_seg = image_tensor.squeeze(0), seg_tensor.squeeze(0)
            registered_seg = (registered_seg > 0.5).int()
        image_with_contours = draw_contours(registered_image, registered_seg)
        image_without_contours = (registered_image * 255.).int()
        
        # save registered images
        frames_with_contours.append(tensor_to_numpy(image_with_contours))
        frames_without_contours.append(tensor_to_numpy(image_without_contours))

        # compute seg area
        seg_np = tensor_to_numpy((seg_tensor > 0.5).int().squeeze(0))
        seg_np = cv2.resize(seg_np.astype(np.uint8), (int(row.xslo), int(row.yslo)))
        area = np.sum(seg_np) * row.scale_x * row.scale_y
        areas_ai.append(area)
        if 'GA Size Final' in row.keys():
            areas_manual.append(row['GA Size Final'])
        else:
            areas_manual.append(float('nan'))
        timepoints.append(row.exdatetime)

    # save thumbnails for videos
    thumbnail_with_contours = frames_with_contours[0].astype(np.uint8)
    thumbnail_without_contours = frames_without_contours[0].astype(np.uint8)
    Image.fromarray(thumbnail_with_contours).save(os.path.join(save_to, 'metadata', 'thumbnails_wcontours', thumbnails_wcontours_name))
    Image.fromarray(thumbnail_without_contours).save(os.path.join(save_to, 'metadata', 'thumbnails_wocontours', thumbnails_wocontours_name))

    # create videos
    clip_with_contours = ImageSequenceClip(frames_with_contours, fps=3)
    clip_without_contours = ImageSequenceClip(frames_without_contours, fps=3)

    # create buffers and save
    clip_with_contours.write_videofile(os.path.join(save_to, 'metadata', 'videos_wcontours', video_wcontours_name))
    clip_without_contours.write_videofile(os.path.join(save_to, 'metadata', 'videos_wocontours', video_wocontours_name))

    # Plot area vs time
    plt.figure(figsize=(6, 4))
    plt.plot(timepoints, areas_ai, 'ro-', color='red')
    plt.plot(timepoints, areas_manual, 'ro-', color='blue')
    plt.legend(['AI-computed Area', 'Manual'])
    plt.xlabel('Time')
    plt.ylabel('GA Area')
    plt.title('GA Area vs Time')
    plt.savefig(os.path.join(save_to, 'metadata', 'plots', plot_name))

    return None

def prepare_presentation_slide(slide, slide_title, video_with_contours_path, video_without_contours_path, thumbnail_with_contours_path, thumbnail_without_contours_path, plot_path):
    title = slide.shapes.title
    title.text = slide_title
    title.text_frame.paragraphs[0].font.size = Pt(40)

    # Add videos
    video_width = Inches(2.5)
    video_height = Inches(2.5)
    left = Inches(1.0)
    top = Inches(1.75)

    # Add video with contours
    slide.shapes.add_movie(video_with_contours_path, left, top, width=video_width, height=video_height, poster_frame_image=thumbnail_with_contours_path)

    # Add video without contours
    slide.shapes.add_movie(video_without_contours_path, left, top + Inches(0.5) + Inches(2.5), width=video_width, height=video_height, poster_frame_image=thumbnail_without_contours_path)

    # Add plot
    plot_width = Inches(6)
    plot_height = Inches(4.5)
    slide.shapes.add_picture(plot_path, Inches(3.75), Inches(2.0), width=plot_width, height=plot_height)

if __name__ == '__main__':

    # load data
    file = 'results/ga_5519923/results.csv'
    save_to = 'results/ga_5519923/powerpoint/'

    df = pd.read_csv(file)
    df['exdatetime'] = pd.to_datetime(df['exdatetime'])
    
    # df1 = pd.read_csv(file)
    # df1['exdatetime'] = pd.to_datetime(df1['exdatetime'])
    
    # df2 = pd.read_csv('data/GA_progression_modelling/af_ga_cases_wmanualarea.csv')
    # df2['exdatetime'] = pd.to_datetime(df2['exdatetime'])

    # merged_df = pd.merge(df1, df2, how='inner', on=['file_path_coris', 'file_path_seg', 'file_path_vessel', 'ptid', 'exdatetime', 'fileeye', 'modality'])
    # merged_df = merged_df.drop(columns=['scale_x_x', 'scale_y_x', 'scale_z',
    #    'size_x', 'size_z', 'xslo_x', 'yslo_x', 'scanpattern', 'start_x',
    #    'start_y', 'end_x', 'end_y'])
    # merged_df = merged_df.rename({'xslo_y': 'xslo', 'yslo_y': 'yslo',
    #    'scale_x_y': 'scale_x', 'scale_y_y': 'scale_y'}, axis=1)
    # merged_df.to_csv('results_2.csv', index=False)

    # Create PowerPoint presentation
    prs = Presentation()
    slide_layout = prs.slide_layouts[5]
    
    for mrn, mrn_df in df.groupby('ptid'):
        for lat, lat_df in mrn_df.groupby('fileeye'):
            lat_df['exdatetime'] = pd.to_datetime(lat_df['exdatetime'])
            lat_df = lat_df.sort_values(by='exdatetime')
            assert isnan(lat_df.iloc[0]['params'])

            # create videos and plots
            create_videos_and_plots(
                lat_df.reset_index(drop=True), 
                save_to=save_to, 
                video_wcontours_name=f'{mrn}_{lat}_wcontours.mp4', 
                video_wocontours_name=f'{mrn}_{lat}_wocontours.mp4', 
                plot_name=f'{mrn}_{lat}_plot.png',
                thumbnails_wcontours_name=f'{mrn}_{lat}_wcontours.png',
                thumbnails_wocontours_name=f'{mrn}_{lat}_wocontours.png'
                )

            # add slides to presentation
            slide = prs.slides.add_slide(slide_layout)
            prepare_presentation_slide(
                slide,
                f'GA area progression analysis\nMRN: {mrn}, Eye: {lat}',
                os.path.join(save_to, 'metadata', 'videos_wcontours', f'{mrn}_{lat}_wcontours.mp4'),
                os.path.join(save_to, 'metadata', 'videos_wocontours', f'{mrn}_{lat}_wocontours.mp4'),
                os.path.join(save_to, 'metadata', 'thumbnails_wcontours', f'{mrn}_{lat}_wcontours.png'),
                os.path.join(save_to, 'metadata', 'thumbnails_wocontours', f'{mrn}_{lat}_wocontours.png'),
                os.path.join(save_to, 'metadata', 'plots', f'{mrn}_{lat}_plot.png')
                )
            
            print(f'Processed mrn: {mrn}, fileeye: {lat}')

    # Save the presentation
    prs.save(os.path.join(save_to, 'segmentation_analysis.pptx'))

            